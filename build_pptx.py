"""
Script de génération de la présentation PowerPoint — Prospect Intelligence
Exécuter : python3 build_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette SG ────────────────────────────────────────────────────────────────
RED     = RGBColor(0xE3, 0x06, 0x13)
BLACK   = RGBColor(0x1A, 0x1A, 0x1A)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREY    = RGBColor(0xF5, 0xF5, 0xF5)
GREY2   = RGBColor(0xE0, 0xE0, 0xE0)
DARK    = RGBColor(0x12, 0x12, 0x12)
MUTED   = RGBColor(0x6B, 0x6B, 0x6B)

LOGO_PATH = "data/sg-logo.png"
SLIDE_W   = Inches(13.33)
SLIDE_H   = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def rgb_to_tuple(rgb):
    return (rgb.red, rgb.green, rgb.blue)

def add_rect(slide, l, t, w, h, fill_rgb, transparency=0):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    shape.line.color.rgb = fill_rgb
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape

def add_text_box(slide, text, l, t, w, h,
                 font_size=18, bold=False, color=BLACK,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox

def add_logo(slide, x=Inches(11.8), y=Inches(0.18), h=Inches(0.6)):
    try:
        pic = slide.shapes.add_picture(LOGO_PATH, x, y, height=h)
        return pic
    except Exception:
        return None

def red_bar(slide, y=Inches(7.1), h=Inches(0.4)):
    """Bottom red bar."""
    add_rect(slide, 0, y, SLIDE_W, h, RED)

def top_band(slide, h=Inches(0.08)):
    """Thin red top band."""
    add_rect(slide, 0, 0, SLIDE_W, h, RED)

def slide_number(slide, num, total=10):
    add_text_box(slide, f"{num} / {total}",
                 Inches(12.5), Inches(7.15), Inches(0.7), Inches(0.28),
                 font_size=9, color=WHITE, align=PP_ALIGN.RIGHT)

def section_title(slide, title, subtitle=""):
    """Red banner with white title on left."""
    add_rect(slide, 0, Inches(0.75), SLIDE_W, Inches(0.9), RED)
    add_text_box(slide, title,
                 Inches(0.35), Inches(0.8), Inches(9), Inches(0.8),
                 font_size=22, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.35), Inches(1.72), Inches(12), Inches(0.4),
                     font_size=13, color=MUTED, italic=True)

def bullet_box(slide, items, l, t, w, h, font_size=13, indent="  •  "):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{indent}{item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = BLACK
        run.font.name = "Calibri"


# ── Slide factory ─────────────────────────────────────────────────────────────

def make_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(layout)

def base_slide(prs, num, total=10):
    """Add blank slide with shared chrome: top band, bottom bar, logo, page number."""
    slide = blank_slide(prs)
    top_band(slide)
    red_bar(slide)
    add_logo(slide)
    slide_number(slide, num, total)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, GREY)   # grey background
    top_band(slide)   # re-draw on top of bg
    red_bar(slide)
    add_logo(slide)
    slide_number(slide, num, total)
    return slide


# ══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════════

prs = make_prs()
TOTAL = 10

# ─── 1. Cover ─────────────────────────────────────────────────────────────────
s = blank_slide(prs)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BLACK)                          # black bg
add_rect(s, 0, 0, Inches(0.18), SLIDE_H, RED)                       # left red stripe
add_rect(s, 0, Inches(6.9), SLIDE_W, Inches(0.6), RED)              # bottom bar
add_logo(s, x=Inches(0.45), y=Inches(0.3), h=Inches(0.75))

add_text_box(s, "PROSPECT INTELLIGENCE",
             Inches(0.45), Inches(1.7), Inches(9), Inches(1.1),
             font_size=36, bold=True, color=WHITE)
add_text_box(s, "Automated HNWI/UHNWI Prospect Detection",
             Inches(0.45), Inches(2.75), Inches(9), Inches(0.6),
             font_size=18, color=RED, bold=True)
add_text_box(s, "Powered by Claude AI · Resend · Supabase · Streamlit",
             Inches(0.45), Inches(3.4), Inches(10), Inches(0.5),
             font_size=13, color=RGBColor(0xAA, 0xAA, 0xAA))

add_text_box(s, "Project Documentation — 2026",
             Inches(0.45), Inches(6.95), Inches(6), Inches(0.4),
             font_size=10, color=WHITE)
slide_number(s, 1, TOTAL)


# ─── 2. Objectif du projet ────────────────────────────────────────────────────
s = base_slide(prs, 2, TOTAL)
section_title(s, "Objectif du Projet",
              "Pourquoi cet outil a été conçu")

add_text_box(s, "Problème",
             Inches(0.35), Inches(2.0), Inches(4), Inches(0.4),
             font_size=14, bold=True, color=RED)
bullet_box(s, [
    "Les signaux de liquidité (IPO, M&A, levées de fonds) sont dispersés sur des centaines de sources",
    "Les commerciaux passent plusieurs heures/semaine à surveiller manuellement la presse",
    "Les prospects les plus chauds ne sont pas identifiés à temps",
], Inches(0.35), Inches(2.4), Inches(6.1), Inches(2.0), font_size=12)

add_text_box(s, "Solution",
             Inches(6.8), Inches(2.0), Inches(4), Inches(0.4),
             font_size=14, bold=True, color=RED)
bullet_box(s, [
    "Pipeline automatisé : collecte → analyse IA → scoring → newsletter",
    "Chaque matin, les 5 meilleurs prospects HNWI/UHNWI livrés par email",
    "Zéro intervention humaine — entièrement piloté par GitHub Actions",
    "Interface d'administration Streamlit pour le suivi et les tests",
], Inches(6.8), Inches(2.4), Inches(6.1), Inches(2.2), font_size=12)

add_rect(s, Inches(6.55), Inches(2.0), Inches(0.04), Inches(2.8), RED)

add_text_box(s, "Résultat attendu",
             Inches(0.35), Inches(5.0), Inches(12), Inches(0.4),
             font_size=13, bold=True, color=BLACK)
add_text_box(s,
    "Réduire à zéro le temps de veille manuelle · Contacter les prospects "
    "au bon moment · Augmenter le taux de conversion grâce à une approche data-driven",
    Inches(0.35), Inches(5.4), Inches(12.4), Inches(0.9),
    font_size=12, color=MUTED)


# ─── 3. Vue d'ensemble du pipeline ───────────────────────────────────────────
s = base_slide(prs, 3, TOTAL)
section_title(s, "Vue d'Ensemble du Pipeline",
              "5 étapes de la collecte à la boîte mail")

steps = [
    ("1", "Collecte", "RSS · DuckDuckGo\nGoogle News"),
    ("2", "Filtrage", "Détection de\nsignaux de liquidité"),
    ("3", "Extraction IA", "Claude Haiku 4.5\n→ Fiche prospect"),
    ("4", "Scoring", "Score 0–100\nUrgence 0–10"),
    ("5", "Newsletter", "Email HTML\n→ Abonnés actifs"),
]

box_w = Inches(2.1)
box_h = Inches(2.0)
gap   = Inches(0.28)
start_x = Inches(0.35)
y_box = Inches(2.1)

for i, (num, title, desc) in enumerate(steps):
    x = start_x + i * (box_w + gap)
    # white card
    add_rect(s, x, y_box, box_w, box_h, WHITE)
    # red number circle (simulated with small rect)
    add_rect(s, x + Inches(0.05), y_box + Inches(0.08),
             Inches(0.42), Inches(0.42), RED)
    add_text_box(s, num, x + Inches(0.05), y_box + Inches(0.08),
                 Inches(0.42), Inches(0.42),
                 font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, title,
                 x + Inches(0.05), y_box + Inches(0.58),
                 box_w - Inches(0.1), Inches(0.45),
                 font_size=13, bold=True, color=BLACK)
    add_text_box(s, desc,
                 x + Inches(0.05), y_box + Inches(1.05),
                 box_w - Inches(0.1), Inches(0.85),
                 font_size=10, color=MUTED)
    # arrow between steps
    if i < len(steps) - 1:
        ax = x + box_w + Inches(0.06)
        add_text_box(s, "›", ax, y_box + Inches(0.7), Inches(0.18), Inches(0.5),
                     font_size=22, bold=True, color=RED, align=PP_ALIGN.CENTER)

add_text_box(s, "Fréquence : quotidienne (07:00 UTC)  ·  Orchestration : GitHub Actions  ·  Durée d'un run : ~2 min",
             Inches(0.35), Inches(4.4), Inches(12.4), Inches(0.4),
             font_size=11, color=MUTED, italic=True, align=PP_ALIGN.CENTER)


# ─── 4. Sources de données ────────────────────────────────────────────────────
s = base_slide(prs, 4, TOTAL)
section_title(s, "Sources de Données",
              "Collecte multi-sources, filtrée aux 7 derniers jours")

sources = [
    ("RSS Financiers", "Reuters · Forbes · TechCrunch · VentureBeat\nPitchBook · Sifted · EU Startups",
     "~60 articles/run", "Gratuit"),
    ("Google News RSS", "Requêtes ciblées IPO · M&A · Fundraising\nFiltrage natif tbs=qdr:w (7 jours)",
     "~20 articles/run", "Gratuit"),
    ("DuckDuckGo News", "6 requêtes avec l'année courante auto-générée\ntimelimit=week pour la fraîcheur",
     "~5–30 articles/run", "Gratuit"),
    ("Tavily (optionnel)", "API News avec filtre days=7\nRésultats enrichis avec contenu complet",
     "~18 articles/run", "1 000 req/mois"),
]

col_w = Inches(3.1)
y0 = Inches(2.05)
for i, (name, detail, volume, cost) in enumerate(sources):
    x = Inches(0.3) + i * (col_w + Inches(0.12))
    add_rect(s, x, y0, col_w, Inches(2.7), WHITE)
    add_rect(s, x, y0, col_w, Inches(0.06), RED)
    add_text_box(s, name, x + Inches(0.12), y0 + Inches(0.15),
                 col_w - Inches(0.2), Inches(0.45),
                 font_size=13, bold=True, color=BLACK)
    add_text_box(s, detail, x + Inches(0.12), y0 + Inches(0.6),
                 col_w - Inches(0.2), Inches(1.0),
                 font_size=10, color=MUTED)
    add_text_box(s, volume, x + Inches(0.12), y0 + Inches(1.65),
                 col_w - Inches(0.2), Inches(0.35),
                 font_size=10, bold=True, color=BLACK)
    add_rect(s, x + Inches(0.12), y0 + Inches(2.15),
             col_w - Inches(0.55), Inches(0.3), GREY2)
    add_text_box(s, cost, x + Inches(0.14), y0 + Inches(2.17),
                 col_w - Inches(0.3), Inches(0.28),
                 font_size=9, color=RED, bold=True)

add_text_box(s,
    "Toutes les sources sont triées par date de publication (plus récent en premier) avant l'analyse IA.",
    Inches(0.35), Inches(5.1), Inches(12.4), Inches(0.4),
    font_size=11, color=MUTED, italic=True)


# ─── 5. Extraction IA ────────────────────────────────────────────────────────
s = base_slide(prs, 5, TOTAL)
section_title(s, "Extraction par Intelligence Artificielle",
              "Claude Haiku 4.5 — analyse de chaque article")

add_text_box(s, "Modèle utilisé : Claude Haiku 4.5 (Anthropic)",
             Inches(0.35), Inches(1.95), Inches(8), Inches(0.4),
             font_size=13, bold=True, color=BLACK)

add_text_box(s, "Pour chaque article retenu, Claude extrait :",
             Inches(0.35), Inches(2.38), Inches(8), Inches(0.35),
             font_size=12, color=MUTED)

fields_left = [
    "Nom complet du prospect",
    "Poste / Titre actuel",
    "Entreprise concernée",
    "Secteur d'activité",
    "Type d'événement (IPO · M&A · Fundraising · Exit…)",
    "Résumé de l'événement",
]
fields_right = [
    "Montant estimé en USD",
    "Localisation géographique",
    "URL de la source",
    "Date de publication",
    "Pitch commercial prêt à l'emploi (sales_pitch)",
    "Score de confiance (0–100) · Score d'urgence (0–10)",
]

bullet_box(s, fields_left,  Inches(0.35), Inches(2.75), Inches(5.8), Inches(2.2), font_size=11)
bullet_box(s, fields_right, Inches(6.6),  Inches(2.75), Inches(6.1), Inches(2.2), font_size=11)
add_rect(s, Inches(6.38), Inches(2.75), Inches(0.04), Inches(2.1), GREY2)

add_rect(s, Inches(0.35), Inches(5.15), Inches(12.4), Inches(1.0), WHITE)
add_rect(s, Inches(0.35), Inches(5.15), Inches(0.06), Inches(1.0), RED)
add_text_box(s, "Coût IA estimé",
             Inches(0.55), Inches(5.2), Inches(3), Inches(0.35),
             font_size=11, bold=True, color=BLACK)
add_text_box(s,
    "~$0.001 par article analysé (Claude Haiku : $1 / 1M tokens entrée · $5 / 1M tokens sortie)\n"
    "Coût total par run : < $0.05 pour 30 articles — soit moins de $1.50 / mois",
    Inches(0.55), Inches(5.55), Inches(12.1), Inches(0.55),
    font_size=11, color=MUTED)


# ─── 6. Scoring & Classement ─────────────────────────────────────────────────
s = base_slide(prs, 6, TOTAL)
section_title(s, "Scoring & Classement des Prospects",
              "Algorithme de priorisation basé sur 4 dimensions")

criteria = [
    ("Montant de l'événement", "40 pts max",
     "IPO / M&A / levée > $1Md → 40 pts\n$100M–$1Md → 20–35 pts\n< $100M → 0–15 pts",
     RED),
    ("Type d'événement", "25 pts max",
     "IPO = 25 pts · Exit = 22 pts\nM&A = 20 pts · Fundraising = 18 pts\nNomination = 10 pts",
     BLACK),
    ("Score de confiance IA", "20 pts max",
     "Confiance Claude (0–100)\npondérée à 20 %\nFiltre min. : 60 pts",
     RED),
    ("Urgence temporelle", "15 pts max",
     "Score d'urgence extrait par Claude\n(0–10) × 1.5\nFaveur aux événements récents",
     BLACK),
]

box_y = Inches(2.1)
bw = Inches(2.95)
for i, (name, pts, desc, col) in enumerate(criteria):
    x = Inches(0.3) + i * (bw + Inches(0.13))
    add_rect(s, x, box_y, bw, Inches(2.85), WHITE)
    add_rect(s, x, box_y, bw, Inches(0.5), col)
    add_text_box(s, name, x + Inches(0.1), box_y + Inches(0.06),
                 bw - Inches(0.15), Inches(0.38),
                 font_size=11, bold=True, color=WHITE)
    add_text_box(s, pts, x + Inches(0.1), box_y + Inches(0.58),
                 bw - Inches(0.15), Inches(0.35),
                 font_size=20, bold=True, color=col)
    add_text_box(s, desc, x + Inches(0.1), box_y + Inches(1.0),
                 bw - Inches(0.15), Inches(1.7),
                 font_size=10, color=MUTED)

add_text_box(s, "Score final (0–100)  =  Σ des 4 critères  ·  Seuil minimum : 40 pts  ·  Top 5 envoyés chaque matin",
             Inches(0.35), Inches(5.2), Inches(12.4), Inches(0.4),
             font_size=12, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

add_rect(s, Inches(0.35), Inches(5.6), Inches(12.4), Inches(0.5), RED)
add_text_box(s,
    "Exemple : Vidit Aatrey (Meesho IPO $1B+) → score 95/100 · urgence 9/10  →  Prospect #1 du jour",
    Inches(0.45), Inches(5.65), Inches(12.0), Inches(0.4),
    font_size=11, color=WHITE)


# ─── 7. La Newsletter ─────────────────────────────────────────────────────────
s = base_slide(prs, 7, TOTAL)
section_title(s, "La Newsletter Quotidienne",
              "Format HTML professionnel · Envoyé via Resend")

add_text_box(s, "Contenu de chaque email",
             Inches(0.35), Inches(2.0), Inches(5.5), Inches(0.38),
             font_size=13, bold=True, color=RED)
bullet_box(s, [
    "En-tête avec logo et date du jour",
    "Fiche détaillée pour chacun des 5 prospects :",
    "     · Nom, poste, entreprise, montant",
    "     · Badge de type d'événement coloré",
    "     · Barre de score de potentiel (0–100)",
    "     · Indicateur d'urgence (●●●● échelle 10)",
    "     · Sales Insight : pitch commercial prêt à l'emploi",
    "     · Lien vers la source originale",
    "Graphique comparatif des 5 prospects (PNG embarqué)",
    "Lien de désinscription → app Streamlit",
], Inches(0.35), Inches(2.4), Inches(6.0), Inches(3.8), font_size=11)

add_text_box(s, "Paramètres d'envoi",
             Inches(7.0), Inches(2.0), Inches(5.5), Inches(0.38),
             font_size=13, bold=True, color=RED)
bullet_box(s, [
    "Plateforme d'envoi : Resend (3 000 emails/mois gratuits)",
    "Format : HTML responsive, compatible Gmail · Outlook",
    "Logo embarqué en base64 (pas de dépendance externe)",
    "Objet : « Top 5 Prospects HNWI — [date] »",
    "Expéditeur : configurable via FROM_EMAIL / FROM_NAME",
    "Destinataires : abonnés actifs uniquement (Supabase)",
    "Désinscription en un clic depuis l'email",
], Inches(7.0), Inches(2.4), Inches(5.9), Inches(3.0), font_size=11)

add_rect(s, Inches(6.75), Inches(2.0), Inches(0.04), Inches(3.8), GREY2)


# ─── 8. Architecture Technique ───────────────────────────────────────────────
s = base_slide(prs, 8, TOTAL)
section_title(s, "Architecture Technique",
              "Stack 100 % cloud · Coût mensuel < $2")

stack = [
    ("Frontend", "Streamlit Community Cloud", "Interface utilisateur · Inscription · Dashboard · Admin", "Gratuit"),
    ("Backend IA", "Claude Haiku 4.5 (Anthropic)", "Extraction prospects · Scoring · Pitch commercial", "< $1.50 / mois"),
    ("Base de données", "Supabase (PostgreSQL)", "Abonnés · Prospects · Logs de run", "Gratuit (500 MB)"),
    ("Envoi email", "Resend API", "Newsletter HTML · Emails de test", "Gratuit (3 000/mois)"),
    ("Automatisation", "GitHub Actions (cron)", "Run quotidien 07:00 UTC · Déclenchement manuel", "Gratuit"),
    ("Collecte news", "RSS + DuckDuckGo + Tavily", "Sources publiques gratuites · 60–100 articles/run", "Gratuit"),
]

row_h = Inches(0.62)
y0 = Inches(2.0)
for i, (layer, tech, desc, cost) in enumerate(stack):
    y = y0 + i * (row_h + Inches(0.05))
    bg = WHITE if i % 2 == 0 else GREY
    add_rect(s, Inches(0.3), y, Inches(12.5), row_h, bg)
    add_text_box(s, layer,  Inches(0.4),  y + Inches(0.12), Inches(1.8), Inches(0.38),
                 font_size=11, bold=True, color=RED)
    add_text_box(s, tech,   Inches(2.3),  y + Inches(0.12), Inches(3.2), Inches(0.38),
                 font_size=11, bold=True, color=BLACK)
    add_text_box(s, desc,   Inches(5.6),  y + Inches(0.12), Inches(5.5), Inches(0.38),
                 font_size=10, color=MUTED)
    add_text_box(s, cost,   Inches(11.2), y + Inches(0.12), Inches(1.5), Inches(0.38),
                 font_size=10, bold=True, color=RED, align=PP_ALIGN.RIGHT)

add_text_box(s, "Coût total estimé : < $2 / mois en production",
             Inches(0.35), Inches(6.7), Inches(12.0), Inches(0.35),
             font_size=12, bold=True, color=BLACK, align=PP_ALIGN.CENTER)


# ─── 9. Déploiement & Automatisation ─────────────────────────────────────────
s = base_slide(prs, 9, TOTAL)
section_title(s, "Déploiement & Automatisation",
              "Mise en ligne en 4 étapes · Zéro maintenance")

steps_deploy = [
    ("1", "Supabase", "Créer le projet gratuit · Exécuter le SQL de création des 3 tables\n(subscribers · prospects · run_log)"),
    ("2", "GitHub Secrets", "Ajouter 6 secrets dans Settings → Actions :\nANTHROPIC_API_KEY · RESEND_API_KEY · FROM_EMAIL\nFROM_NAME · SUPABASE_URL · SUPABASE_KEY"),
    ("3", "Streamlit Cloud", "Connecter le repo GitHub · Renseigner les mêmes\nvariables dans Advanced Settings → Secrets"),
    ("4", "C'est en ligne !", "Le workflow GitHub Actions se déclenche chaque matin\nà 07:00 UTC · L'app est accessible à l'URL publique"),
]

y0 = Inches(2.05)
for i, (num, title, desc) in enumerate(steps_deploy):
    y = y0 + i * Inches(1.08)
    add_rect(s, Inches(0.3), y, Inches(0.55), Inches(0.55), RED)
    add_text_box(s, num, Inches(0.3), y + Inches(0.03),
                 Inches(0.55), Inches(0.5),
                 font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, title, Inches(1.0), y + Inches(0.02),
                 Inches(2.5), Inches(0.35),
                 font_size=13, bold=True, color=BLACK)
    add_text_box(s, desc, Inches(3.6), y + Inches(0.02),
                 Inches(9.1), Inches(0.65),
                 font_size=11, color=MUTED)
    if i < len(steps_deploy) - 1:
        add_rect(s, Inches(0.52), y + Inches(0.55), Inches(0.1), Inches(0.53), GREY2)

add_rect(s, Inches(0.3), Inches(6.4), Inches(12.4), Inches(0.45), BLACK)
add_text_box(s,
    "URL de l'app : https://sg-prospect-intelligence.streamlit.app   ·   "
    "Cron : 0 7 * * 1,3  (lundi & mercredi)",
    Inches(0.4), Inches(6.45), Inches(12.0), Inches(0.35),
    font_size=10, color=WHITE, align=PP_ALIGN.CENTER)


# ─── 10. Synthèse & Bénéfices ─────────────────────────────────────────────────
s = base_slide(prs, 10, TOTAL)
section_title(s, "Synthèse & Bénéfices",
              "Ce que cet outil apporte concrètement")

kpis = [
    ("< 2 min", "Durée d'un run\ncollecte → email"),
    ("< $2", "Coût mensuel\ntotal estimé"),
    ("100 %", "Automatisé\nzéro action manuelle"),
    ("5", "Prospects qualifiés\nchaque matin"),
]
kpi_y = Inches(2.1)
kw = Inches(2.9)
for i, (val, label) in enumerate(kpis):
    x = Inches(0.35) + i * (kw + Inches(0.2))
    add_rect(s, x, kpi_y, kw, Inches(1.3), WHITE)
    add_rect(s, x, kpi_y, kw, Inches(0.06), RED)
    add_text_box(s, val, x + Inches(0.1), kpi_y + Inches(0.12),
                 kw - Inches(0.15), Inches(0.65),
                 font_size=28, bold=True, color=RED, align=PP_ALIGN.CENTER)
    add_text_box(s, label, x + Inches(0.1), kpi_y + Inches(0.75),
                 kw - Inches(0.15), Inches(0.45),
                 font_size=10, color=MUTED, align=PP_ALIGN.CENTER)

add_text_box(s, "Pour les équipes commerciales",
             Inches(0.35), Inches(3.65), Inches(5.9), Inches(0.38),
             font_size=13, bold=True, color=BLACK)
bullet_box(s, [
    "Recevoir chaque matin les 5 prospects les plus chauds directement en boîte mail",
    "Chaque fiche inclut un pitch commercial prêt à l'emploi — gain de temps immédiat",
    "Priorisation par score et urgence — plus jamais de prospect raté",
    "Se désabonner en un clic depuis l'email",
], Inches(0.35), Inches(4.05), Inches(5.9), Inches(2.1), font_size=11)

add_rect(s, Inches(6.55), Inches(3.65), Inches(0.04), Inches(2.5), GREY2)

add_text_box(s, "Pour les équipes techniques",
             Inches(6.8), Inches(3.65), Inches(5.9), Inches(0.38),
             font_size=13, bold=True, color=BLACK)
bullet_box(s, [
    "Code Python modulaire : collecte · traitement · DB · email sont découplés",
    "Base 100 % cloud : Supabase · Streamlit · GitHub Actions",
    "Extensible : ajouter Tavily, changer de modèle IA, nouvelles sources RSS",
    "Monitoring via l'onglet Admin et l'historique des runs en temps réel",
], Inches(6.8), Inches(4.05), Inches(5.9), Inches(2.1), font_size=11)


# ── Sauvegarde ────────────────────────────────────────────────────────────────
OUT = "data/Prospect_Intelligence_Documentation.pptx"
prs.save(OUT)
print(f"✅ Présentation générée : {OUT}")
